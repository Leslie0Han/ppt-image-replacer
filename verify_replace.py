"""
验证脚本：检查替换后的PPT图片内容
"""

from pptx import Presentation
from PIL import Image
import io
import os


def verify_replace():
    """验证替换后的PPT图片"""
    print("=" * 60)
    print("验证 PPT 图片替换结果")
    print("=" * 60)

    # 1. 检查原始PPT
    print("\n1. 原始PPT (测试PPT.pptx):")
    prs = Presentation("测试PPT.pptx")
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                blob = shape.image.blob
                print(f"   第{slide_idx}页: {len(blob)} bytes")

    # 2. 检查新PPT
    print("\n2. 新PPT (测试PPT_updated.pptx):")
    new_prs = Presentation("测试PPT_updated.pptx")
    for slide_idx, slide in enumerate(new_prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                blob = shape.image.blob
                print(f"   第{slide_idx}页: {len(blob)} bytes")

                # 尝试打开图片验证
                try:
                    img = Image.open(io.BytesIO(blob))
                    print(f"         格式: {img.format}, 大小: {img.size}")
                except Exception as e:
                    print(f"         错误: {str(e)}")

    # 3. 检查新图片文件
    print("\n3. 新图片文件:")
    new_images_dir = "new_images"
    for file in os.listdir(new_images_dir):
        filepath = os.path.join(new_images_dir, file)
        size = os.path.getsize(filepath)
        print(f"   {file}: {size} bytes")

    # 4. 比较结果
    print("\n4. 比较结果:")
    original_sizes = []
    new_sizes = []

    prs = Presentation("测试PPT.pptx")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                original_sizes.append(len(shape.image.blob))

    new_prs = Presentation("测试PPT_updated.pptx")
    for slide in new_prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                new_sizes.append(len(shape.image.blob))

    print(f"   原始图片大小: {original_sizes}")
    print(f"   新图片大小: {new_sizes}")

    if original_sizes != new_sizes:
        print("   [OK] 图片已成功替换！")
    else:
        print("   [FAIL] 图片大小未变化，可能替换未生效")

    print("\n" + "=" * 60)


if __name__ == "__main__":
    verify_replace()
