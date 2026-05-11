"""
测试脚本：创建一个包含图片的测试PPT
"""

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def create_test_image(filename, color, size=(200, 200)):
    """创建测试图片"""
    img = Image.new('RGB', size, color)
    img.save(filename)
    return filename


def create_test_ppt():
    """创建测试PPT"""
    # 创建测试图片目录
    test_dir = "test_images"
    os.makedirs(test_dir, exist_ok=True)

    # 创建测试图片
    images = [
        create_test_image(f"{test_dir}/户型图A.png", "red"),
        create_test_image(f"{test_dir}/总图.jpg", "blue"),
        create_test_image(f"{test_dir}/效果图01.png", "green"),
    ]

    # 创建PPT
    prs = Presentation()

    # 第一页：户型图
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    slide1.shapes.add_picture(images[0], Inches(1), Inches(1), Inches(4), Inches(3))

    # 第二页：总图
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.shapes.add_picture(images[1], Inches(0.5), Inches(0.5), Inches(5), Inches(4))

    # 第三页：效果图
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.shapes.add_picture(images[2], Inches(2), Inches(1), Inches(4), Inches(3))

    # 保存PPT
    ppt_path = "测试PPT.pptx"
    prs.save(ppt_path)

    print(f"测试PPT已创建: {ppt_path}")
    print(f"测试图片已创建在: {test_dir}/")
    print("\n图片列表:")
    for img in images:
        print(f"  - {img}")

    return ppt_path, test_dir


if __name__ == "__main__":
    create_test_ppt()
