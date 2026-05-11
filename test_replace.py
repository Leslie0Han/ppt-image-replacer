"""
自动化测试脚本：测试PPT图片替换功能
"""

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import shutil


def test_replace():
    """测试图片替换功能"""
    print("=" * 50)
    print("PPT 图片替换功能测试")
    print("=" * 50)

    # 1. 加载测试PPT
    ppt_path = "测试PPT.pptx"
    if not os.path.exists(ppt_path):
        print(f"错误: 找不到测试PPT: {ppt_path}")
        return

    print(f"\n1. 加载PPT: {ppt_path}")
    prs = Presentation(ppt_path)

    # 2. 统计原始图片
    print("\n2. 统计原始图片:")
    original_images = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 图片类型
                original_images.append({
                    'slide_idx': slide_idx,
                    'slide': slide,
                    'shape': shape,
                    'name': shape.name,
                    'size': len(shape.image.blob)
                })
                print(f"   第{slide_idx}页: {shape.name} ({len(shape.image.blob)} bytes)")

    print(f"\n   共找到 {len(original_images)} 张图片")

    # 3. 替换图片
    new_images_dir = "new_images"
    if not os.path.exists(new_images_dir):
        print(f"\n错误: 找不到新图片目录: {new_images_dir}")
        return

    # 获取新图片列表
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
    new_image_files = []
    for file in os.listdir(new_images_dir):
        ext = os.path.splitext(file)[1].lower()
        if ext in image_extensions:
            new_image_files.append(file)
    new_image_files.sort()

    print(f"\n3. 从 {new_images_dir} 替换图片:")
    print(f"   找到 {len(new_image_files)} 张新图片")
    replaced_count = 0

    for idx, img_info in enumerate(original_images):
        # 按顺序匹配（因为PPT中的图片名称可能都是 "Picture 1"）
        if idx < len(new_image_files):
            new_image_path = os.path.join(new_images_dir, new_image_files[idx])
            try:
                # 读取新图片
                with open(new_image_path, 'rb') as f:
                    new_blob = f.read()

                # 使用正确的方法替换图片
                # 获取图片的关系ID（blip_rId是属性，不是方法）
                rId = img_info['shape']._element.blip_rId
                # 通过幻灯片部分的关系获取图片部分并替换
                slide = img_info['slide']
                # 直接使用关系ID来获取图片部分
                if rId in slide.part.rels:
                    rel = slide.part.rels[rId]
                    if 'image' in rel.reltype:
                        rel.target_part._blob = new_blob
                replaced_count += 1
                print(f"   替换: 第{img_info['slide_idx']}页 -> {new_image_files[idx]}")

            except Exception as e:
                print(f"   失败: 第{img_info['slide_idx']}页 - {str(e)}")
        else:
            print(f"   跳过: 第{img_info['slide_idx']}页 (没有对应的新图片)")

    # 4. 保存新PPT
    new_ppt_path = "测试PPT_updated.pptx"
    prs.save(new_ppt_path)

    print(f"\n4. 保存新PPT: {new_ppt_path}")
    print(f"\n5. 测试结果:")
    print(f"   - 原始图片数: {len(original_images)}")
    print(f"   - 成功替换: {replaced_count}")
    print(f"   - 新文件: {new_ppt_path}")

    # 5. 验证新PPT
    print(f"\n6. 验证新PPT:")
    new_prs = Presentation(new_ppt_path)
    for slide_idx, slide in enumerate(new_prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                print(f"   第{slide_idx}页: {shape.name} ({len(shape.image.blob)} bytes)")

    print("\n" + "=" * 50)
    print("测试完成!")
    print("=" * 50)


if __name__ == "__main__":
    test_replace()
