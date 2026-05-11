"""
PPT 图片批量替换工具 - 可视化版本
解决痛点：每轮方案修改后需要重新导入大量图片到PPT中
功能：可视化显示每页图片，支持单张替换、多选替换和批量替换
"""

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from PIL import Image, ImageTk
import io


class ImageCard:
    """单张图片卡片"""

    def __init__(self, parent, image_info, on_select_callback, on_check_callback, use_grid=False):
        self.parent = parent
        self.image_info = image_info
        self.on_select_callback = on_select_callback
        self.on_check_callback = on_check_callback
        self.new_image_path = None
        self.is_selected = False  # 是否被勾选
        self.use_grid = use_grid

        self.create_ui()

    def create_ui(self):
        """创建图片卡片UI"""
        # 主框架 - 固定大小
        self.frame = ttk.Frame(self.parent, relief="solid", borderwidth=2, width=200, height=320)
        # 不在这里pack或grid，由外部控制布局
        self.frame.pack_propagate(False)

        # 顶部勾选框
        self.check_var = tk.BooleanVar(value=False)
        self.check_btn = ttk.Checkbutton(self.frame, text="勾选替换",
                                         variable=self.check_var,
                                         command=self.on_check)
        self.check_btn.pack(pady=(5, 0))

        # 图片预览区域
        self.preview_frame = ttk.Frame(self.frame, width=180, height=180)
        self.preview_frame.pack(padx=10, pady=5)
        self.preview_frame.pack_propagate(False)

        # 显示原图缩略图
        self.show_thumbnail(self.image_info['thumbnail'])

        # 图片名称
        name_label = ttk.Label(self.frame, text=self.image_info['name'],
                              font=("微软雅黑", 10), wraplength=180, anchor="center")
        name_label.pack(pady=(0, 5))

        # 状态标签
        self.status_var = tk.StringVar(value="原图")
        self.status_label = ttk.Label(self.frame, textvariable=self.status_var,
                                      font=("微软雅黑", 11, "bold"), foreground="gray")
        self.status_label.pack(pady=(0, 5))

        # 选择按钮
        self.select_btn = ttk.Button(self.frame, text="选择新图片",
                                    command=self.select_image, width=18)
        self.select_btn.pack(pady=(0, 10))

    def on_check(self):
        """勾选状态改变"""
        self.is_selected = self.check_var.get()
        if self.on_check_callback:
            self.on_check_callback(self.image_info['id'], self.is_selected)

    def set_checked(self, checked):
        """设置勾选状态"""
        self.check_var.set(checked)
        self.is_selected = checked

    def show_thumbnail(self, image):
        """显示缩略图"""
        for widget in self.preview_frame.winfo_children():
            widget.destroy()

        img = image.copy()
        img.thumbnail((170, 170), Image.Resampling.LANCZOS)
        self.tk_image = ImageTk.PhotoImage(img)

        label = ttk.Label(self.preview_frame, image=self.tk_image)
        label.pack(expand=True)

    def select_image(self):
        """选择新图片"""
        file_path = filedialog.askopenfilename(
            title="选择新图片",
            filetypes=[
                ("图片文件", "*.png *.jpg *.jpeg *.gif *.bmp *.tiff"),
                ("所有文件", "*.*")
            ]
        )
        if file_path:
            self.set_new_image(file_path)

    def set_new_image(self, file_path):
        """设置新图片"""
        self.new_image_path = file_path

        try:
            new_img = Image.open(file_path)
            self.show_thumbnail(new_img)

            filename = os.path.basename(file_path)
            self.status_var.set("已替换")
            self.status_label.configure(foreground="green")

            if self.on_select_callback:
                self.on_select_callback(self.image_info['id'], file_path)

        except Exception as e:
            messagebox.showerror("错误", f"无法加载图片: {str(e)}")

    def get_new_image_path(self):
        """获取新图片路径"""
        return self.new_image_path


class SlideCard:
    """幻灯片卡片"""

    def __init__(self, parent, slide_idx, images, on_select_callback, on_check_callback):
        self.parent = parent
        self.slide_idx = slide_idx
        self.images = images
        self.on_select_callback = on_select_callback
        self.on_check_callback = on_check_callback
        self.image_cards = []

        self.create_ui()

    def create_ui(self):
        """创建幻灯片卡片UI"""
        self.frame = ttk.LabelFrame(self.parent, text=f"第 {self.slide_idx} 页",
                                   padding=10)
        self.frame.pack(fill=tk.X, padx=10, pady=5)

        # 顶部信息栏
        info_frame = ttk.Frame(self.frame)
        info_frame.pack(fill=tk.X, pady=(0, 5))

        count_label = ttk.Label(info_frame,
                               text=f"共 {len(self.images)} 张图片",
                               font=("微软雅黑", 10))
        count_label.pack(side=tk.LEFT)

        # 全选按钮
        self.select_all_var = tk.BooleanVar(value=False)
        select_all_btn = ttk.Checkbutton(info_frame, text="全选本页",
                                         variable=self.select_all_var,
                                         command=self.select_all)
        select_all_btn.pack(side=tk.RIGHT)

        # 图片容器 - 使用网格布局让图片居中
        cards_container = ttk.Frame(self.frame)
        cards_container.pack(fill=tk.X, pady=(0, 10))

        # 让列居中
        cards_container.columnconfigure(tuple(range(len(self.images))), weight=1)

        # 创建每个图片的卡片，使用网格布局居中显示
        for idx, img_info in enumerate(self.images):
            card = ImageCard(cards_container, img_info,
                           self.on_select_callback, self.on_check_callback)
            self.image_cards.append(card)
            # 使用网格布局，让卡片居中
            card.frame.grid(row=0, column=idx, padx=10, pady=10, sticky="n")

    def select_all(self):
        """全选/取消全选本页"""
        checked = self.select_all_var.get()
        for card in self.image_cards:
            card.set_checked(checked)
            if self.on_check_callback:
                self.on_check_callback(card.image_info['id'], checked)

    def get_image_cards(self):
        """获取所有图片卡片"""
        return self.image_cards


class PPTImageReplacer:
    """PPT图片替换工具主类"""

    def __init__(self, root):
        self.root = root
        self.root.title("PPT 图片替换工具 - 可视化版")
        self.root.geometry("1400x900")

        # 数据存储
        self.ppt_path = None
        self.prs = None
        self.slides_data = []
        self.new_images = {}
        self.checked_images = set()  # 勾选的图片ID集合
        self.slide_cards = []

        self.setup_ui()

    def setup_ui(self):
        """设置UI界面"""
        main_frame = ttk.Frame(self.root)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # === 顶部工具栏 ===
        toolbar = ttk.Frame(main_frame)
        toolbar.pack(fill=tk.X, padx=10, pady=10)

        # 左侧：PPT选择
        left_frame = ttk.Frame(toolbar)
        left_frame.pack(side=tk.LEFT)

        ttk.Label(left_frame, text="PPT文件:", font=("微软雅黑", 10)).pack(side=tk.LEFT)
        self.ppt_path_var = tk.StringVar(value="未选择")
        ttk.Label(left_frame, textvariable=self.ppt_path_var,
                 width=25, relief="sunken").pack(side=tk.LEFT, padx=5)
        ttk.Button(left_frame, text="选择PPT", command=self.select_ppt).pack(side=tk.LEFT)
        ttk.Button(left_frame, text="扫描", command=self.scan_ppt).pack(side=tk.LEFT, padx=5)

        # 中间：文件夹选择
        mid_frame = ttk.Frame(toolbar)
        mid_frame.pack(side=tk.LEFT, padx=20)

        ttk.Label(mid_frame, text="图片文件夹:", font=("微软雅黑", 10)).pack(side=tk.LEFT)
        self.folder_path_var = tk.StringVar(value="未选择")
        ttk.Label(mid_frame, textvariable=self.folder_path_var,
                 width=25, relief="sunken").pack(side=tk.LEFT, padx=5)
        ttk.Button(mid_frame, text="选择文件夹", command=self.select_folder).pack(side=tk.LEFT)

        # 右侧：操作按钮
        right_frame = ttk.Frame(toolbar)
        right_frame.pack(side=tk.RIGHT)

        ttk.Button(right_frame, text="导出原图",
                  command=self.export_images).pack(side=tk.LEFT, padx=5)
        ttk.Button(right_frame, text="替换勾选的图片",
                  command=self.replace_checked).pack(side=tk.LEFT, padx=5)
        ttk.Button(right_frame, text="一键替换全部",
                  command=self.replace_all).pack(side=tk.LEFT, padx=5)
        ttk.Button(right_frame, text="保存PPT",
                  command=self.save_ppt).pack(side=tk.LEFT)

        # === 操作说明 ===
        help_frame = ttk.Frame(main_frame)
        help_frame.pack(fill=tk.X, padx=10, pady=(0, 5))

        help_text = "操作说明：① 勾选需要替换的图片（可跨页多选）→ ② 点击「替换勾选的图片」选择一张新图片 → ③ 所有勾选的图片都会被替换为同一张"
        ttk.Label(help_frame, text=help_text, font=("微软雅黑", 9),
                 foreground="blue").pack(side=tk.LEFT)

        # === 分隔线 ===
        ttk.Separator(main_frame, orient=tk.HORIZONTAL).pack(fill=tk.X, padx=10)

        # === 主内容区域（可滚动） ===
        content_frame = ttk.Frame(main_frame)
        content_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        self.canvas = tk.Canvas(content_frame)
        scrollbar_y = ttk.Scrollbar(content_frame, orient=tk.VERTICAL,
                                   command=self.canvas.yview)
        scrollbar_x = ttk.Scrollbar(content_frame, orient=tk.HORIZONTAL,
                                   command=self.canvas.xview)

        self.canvas.configure(yscrollcommand=scrollbar_y.set,
                            xscrollcommand=scrollbar_x.set)

        scrollbar_y.pack(side=tk.RIGHT, fill=tk.Y)
        scrollbar_x.pack(side=tk.BOTTOM, fill=tk.X)
        self.canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        self.content_frame = ttk.Frame(self.canvas)
        self.canvas.create_window((0, 0), window=self.content_frame, anchor=tk.NW)

        self.canvas.bind_all("<MouseWheel>", self.on_mousewheel)

        # === 底部状态栏 ===
        status_frame = ttk.Frame(main_frame)
        status_frame.pack(fill=tk.X, padx=10, pady=5)

        self.status_var = tk.StringVar(value="就绪 - 请选择PPT文件并扫描")
        status_bar = ttk.Label(status_frame, textvariable=self.status_var,
                              relief="sunken", padding=5, font=("微软雅黑", 10))
        status_bar.pack(side=tk.LEFT, fill=tk.X, expand=True)

        self.checked_count_var = tk.StringVar(value="已勾选: 0 张")
        checked_label = ttk.Label(status_frame, textvariable=self.checked_count_var,
                                 relief="sunken", padding=5, font=("微软雅黑", 10, "bold"),
                                 foreground="blue")
        checked_label.pack(side=tk.RIGHT, padx=(10, 0))

    def on_mousewheel(self, event):
        """鼠标滚轮事件"""
        self.canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

    def select_ppt(self):
        """选择PPT文件"""
        file_path = filedialog.askopenfilename(
            title="选择PPT文件",
            filetypes=[("PowerPoint文件", "*.pptx"), ("所有文件", "*.*")]
        )
        if file_path:
            self.ppt_path = file_path
            self.ppt_path_var.set(os.path.basename(file_path))
            self.status_var.set(f"已选择PPT: {os.path.basename(file_path)}")

    def select_folder(self):
        """选择图片文件夹"""
        folder_path = filedialog.askdirectory(title="选择图片文件夹")
        if folder_path:
            self.folder_path = folder_path
            self.folder_path_var.set(os.path.basename(folder_path))
            self.status_var.set(f"已选择图片文件夹: {folder_path}")

            if self.slides_data:
                self.match_folder_images()

    def extract_images_from_shapes(self, shapes, slide, slide_idx, global_image_idx):
        """递归提取形状中的图片（包括组合形状中的图片）"""
        slide_images = []

        for shape in shapes:
            # 如果是组合形状，递归遍历
            if shape.shape_type == 6:  # GROUP
                group_images, global_image_idx = self.extract_images_from_shapes(
                    shape.shapes, slide, slide_idx, global_image_idx
                )
                slide_images.extend(group_images)
            # 如果是图片形状
            elif shape.shape_type == 13:  # PICTURE
                try:
                    image_blob = shape.image.blob
                    image_stream = io.BytesIO(image_blob)
                    pil_image = Image.open(image_stream)

                    image_name = shape.name
                    if not image_name or image_name.startswith("Picture"):
                        image_name = f"图片{global_image_idx + 1}"

                    img_info = {
                        'id': global_image_idx,
                        'slide_idx': slide_idx,
                        'slide': slide,
                        'shape': shape,
                        'name': image_name,
                        'thumbnail': pil_image,
                        'blob': image_blob,
                        'position': (shape.left, shape.top),
                        'size': (shape.width, shape.height)
                    }
                    slide_images.append(img_info)
                    global_image_idx += 1
                except Exception as e:
                    print(f"提取图片时出错: {e}")

        return slide_images, global_image_idx

    def scan_ppt(self):
        """扫描PPT文件"""
        if not self.ppt_path:
            messagebox.showwarning("警告", "请先选择PPT文件！")
            return

        self.status_var.set("正在扫描PPT...")
        self.root.update()

        try:
            self.slides_data = []
            self.new_images = {}
            self.checked_images = set()
            self.slide_cards = []

            for widget in self.content_frame.winfo_children():
                widget.destroy()

            self.prs = Presentation(self.ppt_path)

            global_image_idx = 0
            for slide_idx, slide in enumerate(self.prs.slides, 1):
                # 递归提取所有图片（包括组合形状中的）
                slide_images, global_image_idx = self.extract_images_from_shapes(
                    slide.shapes, slide, slide_idx, global_image_idx
                )

                if slide_images:
                    self.slides_data.append({
                        'slide_idx': slide_idx,
                        'images': slide_images
                    })

                    card = SlideCard(self.content_frame, slide_idx,
                                   slide_images, self.on_image_select, self.on_image_check)
                    self.slide_cards.append(card)

            self.content_frame.update_idletasks()
            self.canvas.configure(scrollregion=self.canvas.bbox("all"))

            total_images = sum(len(s['images']) for s in self.slides_data)
            self.status_var.set(f"扫描完成: {len(self.slides_data)} 页, 共 {total_images} 张图片")

            if hasattr(self, 'folder_path'):
                self.match_folder_images()

        except Exception as e:
            messagebox.showerror("错误", f"扫描PPT时出错: {str(e)}")
            self.status_var.set("扫描失败")

    def on_image_select(self, image_id, new_path):
        """图片选择回调"""
        self.new_images[image_id] = new_path
        self.update_status()

    def on_image_check(self, image_id, checked):
        """图片勾选回调"""
        if checked:
            self.checked_images.add(image_id)
        else:
            self.checked_images.discard(image_id)
        self.checked_count_var.set(f"已勾选: {len(self.checked_images)} 张")

    def match_folder_images(self):
        """匹配文件夹中的图片"""
        if not self.folder_path:
            return

        image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
        folder_images = []
        for file in sorted(os.listdir(self.folder_path)):
            ext = os.path.splitext(file)[1].lower()
            if ext in image_extensions:
                folder_images.append(os.path.join(self.folder_path, file))

        if not folder_images:
            return

        image_idx = 0
        for slide_card in self.slide_cards:
            for img_card in slide_card.get_image_cards():
                if image_idx < len(folder_images):
                    img_card.set_new_image(folder_images[image_idx])
                    self.new_images[img_card.image_info['id']] = folder_images[image_idx]
                    image_idx += 1

        self.update_status()

    def replace_checked(self):
        """替换所有勾选的图片"""
        if not self.checked_images:
            messagebox.showwarning("警告", "请先勾选需要替换的图片！\n\n点击图片卡片上的「勾选替换」复选框来选择图片。")
            return

        # 选择一张新图片
        file_path = filedialog.askopenfilename(
            title="选择新图片（将替换所有勾选的图片）",
            filetypes=[
                ("图片文件", "*.png *.jpg *.jpeg *.gif *.bmp *.tiff"),
                ("所有文件", "*.*")
            ]
        )

        if not file_path:
            return

        # 确认替换
        if not messagebox.askyesno("确认", f"确定要用这张图片替换 {len(self.checked_images)} 张勾选的图片吗？"):
            return

        # 替换所有勾选的图片
        try:
            new_img = Image.open(file_path)

            for slide_card in self.slide_cards:
                for img_card in slide_card.get_image_cards():
                    if img_card.image_info['id'] in self.checked_images:
                        img_card.set_new_image(file_path)
                        self.new_images[img_card.image_info['id']] = file_path

            self.status_var.set(f"已用新图片替换 {len(self.checked_images)} 张勾选的图片")
            messagebox.showinfo("成功", f"已替换 {len(self.checked_images)} 张图片！\n\n请点击「保存PPT」生成新文件。")

        except Exception as e:
            messagebox.showerror("错误", f"无法加载图片: {str(e)}")

    def replace_all(self):
        """一键替换所有图片"""
        if not self.new_images:
            messagebox.showwarning("警告", "没有可替换的图片！请先选择新图片。")
            return

        if not messagebox.askyesno("确认", f"确定要替换 {len(self.new_images)} 张图片吗？"):
            return

        self.save_ppt()

    def export_images(self):
        """导出PPT中的所有图片"""
        if not self.slides_data:
            messagebox.showwarning("警告", "请先扫描PPT文件！")
            return

        # 选择导出目录
        export_dir = filedialog.askdirectory(title="选择导出目录")
        if not export_dir:
            return

        self.status_var.set("正在导出图片...")
        self.root.update()

        try:
            exported_count = 0

            for slide_data in self.slides_data:
                slide_idx = slide_data['slide_idx']

                for img_idx, img_info in enumerate(slide_data['images'], 1):
                    # 获取图片格式
                    image_blob = img_info['blob']
                    image_stream = io.BytesIO(image_blob)
                    pil_image = Image.open(image_stream)

                    # 确定文件扩展名
                    format_name = pil_image.format
                    if format_name == 'JPEG':
                        ext = '.jpg'
                    elif format_name == 'PNG':
                        ext = '.png'
                    elif format_name == 'GIF':
                        ext = '.gif'
                    elif format_name == 'BMP':
                        ext = '.bmp'
                    else:
                        ext = '.png'  # 默认使用png

                    # 生成文件名：第几页第几张图
                    filename = f"第{slide_idx}页第{img_idx}张图{ext}"
                    filepath = os.path.join(export_dir, filename)

                    # 保存图片
                    pil_image.save(filepath)
                    exported_count += 1

            self.status_var.set(f"导出完成: 共导出 {exported_count} 张图片")
            messagebox.showinfo("成功",
                              f"图片导出完成！\n\n"
                              f"已导出: {exported_count} 张图片\n"
                              f"保存位置: {export_dir}")

            # 打开导出目录
            os.startfile(export_dir)

        except Exception as e:
            messagebox.showerror("错误", f"导出图片时出错: {str(e)}")
            self.status_var.set("导出失败")

    def save_ppt(self):
        """保存PPT"""
        if not self.prs:
            messagebox.showwarning("警告", "请先扫描PPT文件！")
            return

        if not self.new_images:
            messagebox.showwarning("警告", "没有可替换的图片！")
            return

        self.status_var.set("正在保存PPT...")
        self.root.update()

        try:
            replaced_count = 0
            for image_id, new_path in self.new_images.items():
                for slide_data in self.slides_data:
                    for img_info in slide_data['images']:
                        if img_info['id'] == image_id:
                            with open(new_path, 'rb') as f:
                                new_blob = f.read()

                            rId = img_info['shape']._element.blip_rId
                            slide = img_info['slide']
                            if rId in slide.part.rels:
                                rel = slide.part.rels[rId]
                                if 'image' in rel.reltype:
                                    rel.target_part._blob = new_blob
                                    replaced_count += 1
                            break

            ppt_dir = os.path.dirname(self.ppt_path)
            ppt_name = os.path.splitext(os.path.basename(self.ppt_path))[0]
            new_ppt_path = os.path.join(ppt_dir, f"{ppt_name}_updated.pptx")

            self.prs.save(new_ppt_path)

            self.status_var.set(f"保存完成: 已替换 {replaced_count} 张图片")
            messagebox.showinfo("成功",
                              f"PPT已保存！\n\n"
                              f"已替换: {replaced_count} 张图片\n"
                              f"新文件: {new_ppt_path}")

            os.startfile(ppt_dir)

        except Exception as e:
            messagebox.showerror("错误", f"保存PPT时出错: {str(e)}")
            self.status_var.set("保存失败")

    def update_status(self):
        """更新状态栏"""
        total_images = sum(len(s['images']) for s in self.slides_data)
        replaced_count = len(self.new_images)
        self.status_var.set(f"已扫描 {len(self.slides_data)} 页, "
                          f"共 {total_images} 张图片, "
                          f"已选择替换 {replaced_count} 张")


def main():
    root = tk.Tk()
    app = PPTImageReplacer(root)
    root.mainloop()


if __name__ == "__main__":
    main()
